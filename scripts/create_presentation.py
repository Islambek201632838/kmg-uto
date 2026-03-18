"""Генерация презентации ИС УТО для хакатона."""
import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCREENSHOTS = "front_screenshots"
OUT = "presentation_ИС_УТО.pptx"

# Цвета
BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xFF, 0x2D, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
GREEN = RGBColor(0x00, 0xC8, 0x53)
YELLOW = RGBColor(0xFF, 0xD6, 0x00)
BLUE = RGBColor(0x21, 0x96, 0xF3)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_bullet(tf, text, size=16, color=LIGHT, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(4)

def add_image(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        kwargs = {"left": Inches(left), "top": Inches(top)}
        if width: kwargs["width"] = Inches(width)
        if height: kwargs["height"] = Inches(height)
        slide.shapes.add_picture(path, **kwargs)

def add_metric_box(slide, left, top, label, value, color=GREEN):
    # Фон
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.5), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x40)
    shape.line.fill.background()
    # Метка
    add_text(slide, left + 0.15, top + 0.1, 2.2, 0.4, label, size=12, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, left + 0.15, top + 0.45, 2.2, 0.6, value, size=28, bold=True, color=color, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 1: Титульный
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1, 1.5, 11, 1.5, "ИС УТО", size=60, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.0, 11, 1, "Интеллектуальная система маршрутизации спецтехники", size=28, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 1, 4.2, 11, 0.8, "Месторождение Жетыбай  |  KMG Hackathon 2026", size=18, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(slide, 1, 5.5, 11, 0.8, "Vehicle Routing Problem (VRP) с временными окнами, multi-depot, open-end", size=16, color=LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 2: Проблема
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "Проблема", size=36, bold=True, color=ACCENT)

tf = add_text(slide, 0.5, 1.2, 6, 5, "", size=16)
tf.paragraphs[0].text = ""
items = [
    ("Ручное назначение", "Диспетчер вручную подбирает технику на глаз"),
    ("Нет расчёта ETA", "Маршрут и время прибытия не рассчитываются"),
    ("Холостой пробег", "Техника проезжает лишние километры"),
    ("Нет multi-stop", "Возможность объединить заявки не используется"),
    ("Нет приоритетов", "Срочные заявки обрабатываются наравне с обычными"),
]
for title, desc in items:
    add_bullet(tf, f"❌  {title}: {desc}", size=16, color=LIGHT)

add_text(slide, 7, 1.5, 5.5, 1.5, "52 единицы техники\n3 450 скважин\nРучное планирование", size=22, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
add_text(slide, 7, 3.5, 5.5, 2, "Последствия:\n• Рост холостого пробега\n• Задержки срочных работ\n• Неравномерная загрузка парка\n• Снижение эффективности", size=16, color=LIGHT)


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 3: Решение
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "Решение: ИС УТО", size=36, bold=True, color=ACCENT)

tf = add_text(slide, 0.5, 1.3, 6, 5, "", size=16)
tf.paragraphs[0].text = ""
solutions = [
    "✅  Рекомендация топ-3 техники с обоснованием (score + reason)",
    "✅  Маршруты по графу дорог (Dijkstra, 4624 узла, 38062 ребра)",
    "✅  Multi-stop группировка заявок (OR-Tools CVRPTW)",
    "✅  Сравнение baseline vs оптимизированный (79% экономия)",
    "✅  LLM обоснование через Gemini 2.0 Flash",
    "✅  Визуализация на карте (Folium + Streamlit)",
    "✅  REST API (FastAPI) + Swagger документация",
]
for s in solutions:
    add_bullet(tf, s, size=15, color=LIGHT)

add_text(slide, 7, 1.5, 5.5, 1, "Результат:", size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_metric_box(slide, 7.2, 2.3, "Baseline", "3 624 км", YELLOW)
add_metric_box(slide, 10, 2.3, "VRP", "760 км", GREEN)
add_metric_box(slide, 8.5, 4.0, "Экономия", "79.0%", ACCENT)


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 4: Архитектура
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "Архитектура системы", size=36, bold=True, color=ACCENT)

# Левая колонка — стек
tf = add_text(slide, 0.5, 1.3, 4, 5.5, "", size=14)
tf.paragraphs[0].text = ""
stack = [
    ("Backend", "Python 3.11, FastAPI, async/await"),
    ("БД", "PostgreSQL 17, SQLAlchemy 2.0, Alembic"),
    ("Граф", "networkx DiGraph, KDTree (scipy)"),
    ("VRP", "Google OR-Tools CVRPTW"),
    ("SA", "Simulated Annealing (сравнение)"),
    ("LLM", "Gemini 2.0 Flash (reason)"),
    ("UI", "Streamlit + Folium карта"),
    ("Infra", "Docker Compose (3 сервиса)"),
]
for comp, tech in stack:
    add_bullet(tf, f"{comp}:  {tech}", size=14, color=LIGHT)

# Правая колонка — flow
tf2 = add_text(slide, 5, 1.3, 7.5, 5.5, "", size=14)
tf2.paragraphs[0].text = "Lifespan Init (при старте):"
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.color.rgb = YELLOW
tf2.paragraphs[0].font.size = Pt(16)
steps = [
    "1. Граф дорог → networkx.DiGraph (0.4с)",
    "2. KDTree → snap_to_node за O(log N) (0.003с)",
    "3. 3293 скважины → привязка к узлам (0.01с)",
    "4. 117 единиц техники → позиции, скорости, skills (0.3с)",
    "5. 25 заявок → time windows, SLA (0.001с)",
    "6. Smoke tests всех сервисов",
    "",
    "Concurrency:",
    "• ThreadPoolExecutor(8) — Dijkstra параллельно",
    "• async/await — FastAPI endpoints",
    "• SSSP reverse graph — O(1) vs O(N) lookup",
]
for s in steps:
    add_bullet(tf2, s, size=13, color=LIGHT, bold=("Concurrency" in s))


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 5: Алгоритмы и формулы
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "Алгоритмы и математическая модель", size=36, bold=True, color=ACCENT)

# Целевая функция
add_text(slide, 0.5, 1.2, 12, 0.5, "Целевая функция:", size=18, bold=True, color=YELLOW)
add_text(slide, 0.5, 1.7, 12, 0.5, "min Z = α·Σ(пробег) + β·Σ(wₗ·τₗ) + γ·Σ(простой)", size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Скоринг
add_text(slide, 0.5, 2.5, 6, 0.5, "Формула скоринга (из PPTX формализации):", size=16, bold=True, color=YELLOW)
add_text(slide, 0.5, 3.0, 6, 0.4, "score = 1 - (0.30·D + 0.30·ETA + 0.15·wait + 0.25·SLA)", size=15, color=WHITE)

tf = add_text(slide, 0.5, 3.5, 5.5, 3, "", size=13)
tf.paragraphs[0].text = ""
add_bullet(tf, "ωd = 0.30 — расстояние по графу", size=13, color=LIGHT)
add_bullet(tf, "ωt = 0.30 — ETA (travel + wait)", size=13, color=LIGHT)
add_bullet(tf, "ωw = 0.15 — время ожидания (free_at)", size=13, color=LIGHT)
add_bullet(tf, "ωp = 0.25 — штраф SLA (high +2ч / med +5ч / low +12ч)", size=13, color=LIGHT)

# Ограничения
add_text(slide, 7, 2.5, 5.5, 0.5, "Ограничения:", size=16, bold=True, color=YELLOW)
tf3 = add_text(slide, 7, 3.0, 5.5, 4, "", size=13)
tf3.paragraphs[0].text = ""
constraints = [
    "Hard 1: Каждая заявка = 1 машина",
    "Hard 2: Совместимость техника↔работа",
    "Hard 3: Маршрут по графу дорог",
    "Hard 4: Open-end (не возвращается)",
    "Hard 5: Начало ≥ прибытие",
    "Soft 6: Временные окна → штраф",
    "Soft 7: SLA дедлайн → штраф",
    "Soft 8: Multi-stop detour ≤ 30%",
]
for c in constraints:
    color = GREEN if "Hard" in c else YELLOW
    add_bullet(tf3, c, size=12, color=color)


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 6: Демо — Рекомендации
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "Демо: Рекомендация техники", size=36, bold=True, color=ACCENT)
add_image(slide, f"{SCREENSHOTS}/img.png", 0.3, 1.2, width=6.3)
add_image(slide, f"{SCREENSHOTS}/img_1.png", 6.8, 1.2, width=6.0)
add_text(slide, 0.5, 5.5, 12, 1, "POST /api/recommendations → Топ-3 техники с score, ETA, расстоянием и LLM-обоснованием", size=14, color=LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 7: Демо — Маршрут
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "Демо: Маршрут по графу дорог", size=36, bold=True, color=ACCENT)
add_image(slide, f"{SCREENSHOTS}/img_2.png", 1.5, 1.2, width=10)
add_text(slide, 0.5, 5.8, 12, 1, "POST /api/route → Dijkstra кратчайший путь, polyline на спутниковой карте, 3 слоя", size=14, color=LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 8: Демо — Мультизадачность
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "Демо: Группировка заявок (multi-stop)", size=36, bold=True, color=ACCENT)
add_image(slide, f"{SCREENSHOTS}/img_3.png", 0.3, 1.2, width=6.3)
add_image(slide, f"{SCREENSHOTS}/img_4.png", 6.8, 1.2, width=6.0)
add_text(slide, 0.5, 5.5, 12, 1, "POST /api/multitask → OR-Tools VRP группирует заявки, считает savings vs baseline", size=14, color=LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 9: Сравнение 3 алгоритмов
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "Сравнение: Baseline vs SA vs VRP", size=36, bold=True, color=ACCENT)
add_image(slide, f"{SCREENSHOTS}/img_5.png", 0.3, 1.2, width=6.3)
add_image(slide, f"{SCREENSHOTS}/img_6.png", 6.8, 1.2, width=6.0)

# Таблица результатов
add_metric_box(slide, 0.5, 5.3, "Baseline", "3 624 км", YELLOW)
add_metric_box(slide, 3.5, 5.3, "SA (18K iter)", "1 179 км", BLUE)
add_metric_box(slide, 6.5, 5.3, "VRP (OR-Tools)", "760 км", GREEN)
add_metric_box(slide, 9.5, 5.3, "Экономия VRP", "79.0%", ACCENT)


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 10: Демо-сценарии для жюри
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "Демо-сценарии для жюри", size=36, bold=True, color=ACCENT)
add_image(slide, f"{SCREENSHOTS}/img_7.png", 0.2, 1.1, width=4.2)
add_image(slide, f"{SCREENSHOTS}/img_8.png", 4.5, 1.1, width=4.2)
add_image(slide, f"{SCREENSHOTS}/img_9.png", 8.8, 1.1, width=4.2)

add_text(slide, 0.2, 5.2, 4.2, 1, "Сценарий 1:\nСрочная заявка (high)", size=13, color=YELLOW, align=PP_ALIGN.CENTER)
add_text(slide, 4.5, 5.2, 4.2, 1, "Сценарий 2:\nBaseline vs VRP (79%)", size=13, color=GREEN, align=PP_ALIGN.CENTER)
add_text(slide, 8.8, 5.2, 4.2, 1, "Сценарий 3:\nMulti-stop (33.5%)", size=13, color=BLUE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 11: LLM + Инновации
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "Инновации", size=36, bold=True, color=ACCENT)

add_image(slide, f"{SCREENSHOTS}/img_10.png", 7, 1.2, width=5.8)

add_text(slide, 0.5, 1.3, 6, 0.5, "LLM для обоснований (Gemini 2.0 Flash)", size=18, bold=True, color=YELLOW)
tf = add_text(slide, 0.5, 1.9, 6, 2, "", size=14)
tf.paragraphs[0].text = ""
add_bullet(tf, "Каждая рекомендация → LLM объяснение на русском", size=14, color=LIGHT)
add_bullet(tf, "Анализирует: расстояние, ETA, SLA, совместимость", size=14, color=LIGHT)
add_bullet(tf, "Fallback на шаблон если API недоступен", size=14, color=LIGHT)

add_text(slide, 0.5, 4.0, 6, 0.5, "3 подхода к оптимизации", size=18, bold=True, color=YELLOW)
tf2 = add_text(slide, 0.5, 4.6, 6, 2.5, "", size=14)
tf2.paragraphs[0].text = ""
add_bullet(tf2, "Baseline — жадный, ближайшая свободная", size=14, color=LIGHT)
add_bullet(tf2, "Simulated Annealing — метаэвристика, 100K итераций", size=14, color=LIGHT)
add_bullet(tf2, "OR-Tools CVRPTW — MILP-солвер (лучший результат)", size=14, color=LIGHT)


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 12: Масштабирование
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "Потенциал масштабирования", size=36, bold=True, color=ACCENT)

tf = add_text(slide, 0.5, 1.3, 5.5, 5.5, "", size=15)
tf.paragraphs[0].text = ""
items = [
    ("Другие месторождения", "Загрузить новый граф дорог → готово"),
    ("Больше техники", "1000+ единиц — scipy sparse matrix"),
    ("Реальные заявки", "CSV/API интеграция с КМГ системами"),
    ("Realtime GPS", "Wialon WebSocket → live позиции"),
    ("ML скорости", "Предсказание avg_speed по погоде, типу дороги"),
    ("Kubernetes", "Горизонтальное масштабирование API"),
]
for title, desc in items:
    add_bullet(tf, f"→  {title}", size=16, color=GREEN, bold=True)
    add_bullet(tf, f"     {desc}", size=14, color=LIGHT)

add_text(slide, 7, 1.3, 5.5, 0.5, "Текущие метрики:", size=18, bold=True, color=YELLOW)
add_metric_box(slide, 7, 2.0, "Узлов графа", "4 624", BLUE)
add_metric_box(slide, 10, 2.0, "Рёбер", "38 062", BLUE)
add_metric_box(slide, 7, 3.5, "Техника", "117 ед.", GREEN)
add_metric_box(slide, 10, 3.5, "Скважины", "3 293", GREEN)
add_metric_box(slide, 7, 5.0, "Init time", "< 1 сек", YELLOW)
add_metric_box(slide, 10, 5.0, "API response", "< 8 сек", YELLOW)


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 13: Итоги
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1, 0.8, 11, 1, "Итоги", size=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_metric_box(slide, 0.8, 2.2, "Алгоритм (35%)", "VRP + SA + Baseline", GREEN)
add_metric_box(slide, 3.6, 2.2, "Исполнение (25%)", "FastAPI + Docker", BLUE)
add_metric_box(slide, 6.4, 2.2, "Инновации (20%)", "LLM + 3 подхода", YELLOW)
add_metric_box(slide, 9.2, 2.2, "Применимость (15%)", "3 демо + UI", GREEN)

add_text(slide, 1, 4.0, 11, 1, "79% экономии расстояния на 25 заявках", size=32, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(slide, 1, 5.0, 11, 0.8, "3 624 км (baseline) → 760 км (VRP) — сокращение холостого пробега в 4.8 раза", size=18, color=LIGHT, align=PP_ALIGN.CENTER)

add_text(slide, 1, 6.0, 11, 1, "docker compose up -d  →  всё работает", size=20, color=YELLOW, align=PP_ALIGN.CENTER)
add_text(slide, 1, 6.6, 11, 0.5, "API: localhost:8003/docs  |  UI: localhost:8503", size=16, color=LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Презентация сохранена: {OUT}")
print(f"Слайдов: {len(prs.slides)}")
