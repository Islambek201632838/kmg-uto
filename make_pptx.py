"""ИС УТО — 3-min pitch presentation (10 slides)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK = RGBColor(0x1a, 0x1a, 0x2e)
BLUE = RGBColor(0x1e, 0x3a, 0x8a)
MID = RGBColor(0x2d, 0x6a, 0xd4)
GREEN = RGBColor(0x22, 0xc5, 0x5e)
RED = RGBColor(0xef, 0x44, 0x44)
AMBER = RGBColor(0xf5, 0x9e, 0x0b)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF0, 0xF4, 0xFF)
TEXT = RGBColor(0x1a, 0x1a, 0x2e)
GRAY = RGBColor(0x94, 0xa3, 0xb8)
PURPLE = RGBColor(0x6a, 0x1b, 0x9a)

def rgb(r, g, b): return RGBColor(r, g, b)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 10

def rect(s, l, t, w, h, fill):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    sh.fill.solid(); sh.fill.fore_color.rgb = fill

def txt(s, text, l, t, w, h, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.italic = italic

def header(s, title, sub=""):
    rect(s, 0, 0, 13.33, 1.3, BLUE)
    txt(s, title, 0.4, 0.12, 10, 0.65, size=28, bold=True, color=WHITE)
    if sub:
        txt(s, sub, 0.4, 0.78, 11, 0.4, size=13, color=AMBER, italic=True)

def footer(s, n):
    txt(s, "http://89.207.255.254:8503", 0.3, 7.1, 5, 0.3, size=10, color=GRAY, italic=True)
    txt(s, f"{n} / {TOTAL}", 12.0, 7.1, 1.2, 0.3, size=10, color=GRAY, align=PP_ALIGN.RIGHT)

# ═══ SLIDE 1 — TITLE ═══
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, BLUE)
rect(s, 0, 5.8, 13.33, 0.18, AMBER)
txt(s, "ИС УТО", 0.5, 1.0, 12.33, 1.2, size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Интеллектуальная система маршрутизации спецтехники", 0.5, 2.5, 12.33, 0.7, size=22, color=AMBER, align=PP_ALIGN.CENTER)
txt(s, "OR-Tools VRP  ·  Dijkstra  ·  KDTree  ·  FastAPI  ·  Streamlit  ·  Folium", 0.5, 3.3, 12.33, 0.5, size=14, color=rgb(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)
txt(s, "Хакатон КМГ 2026  ·  Месторождение Жетыбай", 0.5, 5.0, 12.33, 0.5, size=13, color=rgb(0x99,0xAA,0xCC), align=PP_ALIGN.CENTER, italic=True)
footer(s, 1)

# ═══ SLIDE 2 — PROBLEM ═══
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "Проблема: ручная диспетчеризация", "Диспетчер назначает технику вручную — без расчёта маршрута и ETA")
cards = [
    (RED, "Холостой пробег", "Техника едет по ощущениям,\nне по оптимальному маршруту"),
    (AMBER, "Нет ETA", "Диспетчер не знает когда\nтехника доберётся до скважины"),
    (MID, "Нет группировки", "3 заявки рядом обслуживаются\n3 разными машинами"),
    (PURPLE, "Нет приоритетов", "Срочная заявка ждёт наравне\nс плановой"),
]
for i, (color, title, body) in enumerate(cards):
    col, row = i % 2, i // 2
    lx, ty = 0.3 + col * 6.52, 1.5 + row * 2.75
    rect(s, lx, ty, 6.22, 2.5, WHITE)
    rect(s, lx, ty, 0.1, 2.5, color)
    txt(s, title, lx + 0.25, ty + 0.15, 5.7, 0.4, size=15, bold=True, color=color)
    txt(s, body, lx + 0.25, ty + 0.65, 5.7, 1.6, size=13, color=TEXT)
footer(s, 2)

# ═══ SLIDE 3 — SOLUTION (3 steps) ═══
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "Решение: AI-диспетчер за 3 шага")
steps = [
    (BLUE, "Шаг 1", "Заявка → snap_to_node", "Скважина привязывается к ближайшему\nузлу дорожного графа (KDTree, O(log N))"),
    (MID, "Шаг 2", "Dijkstra → маршрут + ETA", "Кратчайший путь по графу дорог\n4 624 узла, 38 062 рёбра"),
    (GREEN, "Шаг 3", "OR-Tools VRP → назначение", "Оптимальная техника с учётом\nприоритета, SLA, skills, time windows"),
]
for i, (color, tag, title, body) in enumerate(steps):
    lx = 0.3 + i * 4.25
    rect(s, lx, 1.5, 4.0, 5.5, WHITE)
    rect(s, lx, 1.5, 4.0, 0.5, color)
    txt(s, tag, lx + 0.15, 1.55, 1.5, 0.4, size=11, bold=True, color=AMBER)
    txt(s, title, lx + 0.15, 2.1, 3.7, 0.5, size=14, bold=True, color=TEXT)
    txt(s, body, lx + 0.15, 2.7, 3.7, 3.5, size=12, color=TEXT)
footer(s, 3)

# ═══ SLIDE 4 — QR FRONTEND ═══
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, BLUE)
rect(s, 0, 6.2, 13.33, 0.18, AMBER)
txt(s, "Попробуйте прямо сейчас", 0.5, 0.3, 12.33, 0.7, size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Откройте камеру телефона и наведите на QR-код", 0.5, 1.0, 12.33, 0.4, size=15, color=AMBER, align=PP_ALIGN.CENTER)
try:
    s.shapes.add_picture("qr_streamlit.png", Inches(4.67), Inches(1.6), width=Inches(4.0), height=Inches(4.0))
except: pass
txt(s, "http://89.207.255.254:8503", 0.5, 5.75, 12.33, 0.4, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Рекомендации · Маршрут · Мультизадачность · Сравнение · Демо", 0.5, 6.3, 12.33, 0.35, size=12, color=rgb(0xAA,0xCC,0xFF), align=PP_ALIGN.CENTER)
footer(s, 4)

# ═══ SLIDE 5 — RESULTS (79% savings) ═══
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "Результат: экономия 79% пробега", "Baseline (greedy) vs OR-Tools VRP vs Simulated Annealing")
metrics = [
    (RED, "3 624 км", "Baseline\n(ближайшая техника)"),
    (AMBER, "1 179 км", "Simulated Annealing\n(67.5% экономия)"),
    (GREEN, "760 км", "OR-Tools VRP\n(79% экономия)"),
]
for i, (color, val, label) in enumerate(metrics):
    lx = 0.3 + i * 4.25
    rect(s, lx, 1.5, 4.0, 2.5, color)
    txt(s, val, lx + 0.1, 1.7, 3.8, 0.9, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, label, lx + 0.1, 2.8, 3.8, 0.9, size=14, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "25 заявок · 117 единиц техники · 4 624 узла дорожного графа", 0.5, 4.3, 12.33, 0.4, size=14, color=TEXT, align=PP_ALIGN.CENTER)

rect(s, 0.3, 4.9, 12.73, 2.3, WHITE)
txt(s, "Скоринг: score = 1 − (0.30·Distance + 0.30·ETA + 0.15·Wait + 0.25·SLA_penalty)", 0.5, 5.0, 12, 0.4, size=12, bold=True, color=BLUE)
txt(s, "VRP: OR-Tools CVRPTW — multi-depot, time windows, open-end, skill compatibility\n"
       "Dijkstra: кратчайший путь по графу (networkx, ThreadPool×8)\n"
       "KDTree: snap_to_node O(log N) — 3 293 скважины привязаны к графу за 30ms", 0.5, 5.5, 12, 1.5, size=11, color=TEXT)
footer(s, 5)

# ═══ SLIDE 6 — API + QR DOCS ═══
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "API — 6 эндпоинтов", "FastAPI · OpenAPI /docs · Pydantic · Docker Compose")
try:
    s.shapes.add_picture("qr_docs.png", Inches(5.17), Inches(1.5), width=Inches(3.0), height=Inches(3.0))
except: pass
txt(s, "http://89.207.255.254:8003/docs", 0.5, 4.6, 12.33, 0.4, size=16, bold=True, color=MID, align=PP_ALIGN.CENTER)
eps = [
    ("POST", "/api/recommendations", "Топ-3 техники + score + reason (LLM)"),
    ("POST", "/api/route", "Маршрут по графу: distance, time, coords"),
    ("POST", "/api/multitask", "Группировка заявок + savings vs baseline"),
    ("GET", "/api/compare", "Baseline vs VRP vs SA — 3 алгоритма"),
    ("GET", "/api/tasks", "25 заявок с координатами"),
    ("GET", "/api/fleet", "117 единиц техники GPS"),
]
for i, (method, path, desc) in enumerate(eps):
    ty = 5.15 + i * 0.35
    color = GREEN if method == "POST" else MID
    txt(s, f"{method}  {path}", 0.5, ty, 5.5, 0.3, size=10, bold=True, color=color)
    txt(s, desc, 6.2, ty, 6.5, 0.3, size=10, color=TEXT)
footer(s, 6)

# ═══ SLIDE 7 — CRITERIA ═══
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "Критерии хакатона — всё покрыто")
criteria = [
    (35, BLUE, "Качество алгоритма", "OR-Tools CVRPTW: 79% экономия. Baseline + SA для сравнения.\nDijkstra по графу, скоринг 4-компонентный"),
    (25, MID, "Техническое исполнение", "FastAPI async, ThreadPool×8, Docker Compose, Pydantic,\nsmoke tests, Alembic, PostgreSQL 17"),
    (20, GREEN, "Инновационность", "Gemini LLM для обоснования выбора, Simulated Annealing,\nSSP reverse-graph O(N) вместо O(N²)"),
    (15, AMBER, "Практическая применимость", "117 машин, 3293 скважины, Streamlit UI 5 режимов,\n3 демо-сценария, Folium карты"),
    (10, PURPLE, "Обработка данных", "KDTree snap 30ms, bulk SQL, parallel Dijkstra,\nVRP 25 задач за 5 сек"),
]
for i, (w, color, title, body) in enumerate(criteria):
    ty = 1.5 + i * 1.1
    rect(s, 0.3, ty, 0.9, 0.9, color)
    txt(s, f"{w}%", 0.3, ty + 0.2, 0.9, 0.5, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 1.35, ty, 11.68, 0.9, WHITE if i % 2 == 0 else rgb(0xF5,0xF8,0xFF))
    txt(s, title, 1.5, ty + 0.05, 4.0, 0.35, size=13, bold=True, color=color)
    txt(s, body, 1.5, ty + 0.42, 11.3, 0.45, size=10, color=TEXT)
footer(s, 7)

# ═══ SLIDE 8 — VRP EXPLAINED ═══
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "VRP: что это и почему OR-Tools", "Vehicle Routing Problem — задача маршрутизации транспорта")
features = [
    (BLUE, "Multi-depot", "Каждая машина стартует из своей GPS-позиции\n(не из одной базы)"),
    (MID, "Time Windows", "Заявки привязаны к сменам:\nдень 08:00–20:00, ночь 20:00–08:00"),
    (GREEN, "Multi-stop", "Одна машина может выполнить несколько\nзаявок за одну смену подряд"),
    (AMBER, "Open-end", "Техника НЕ возвращается на базу —\nостаётся на последней скважине"),
    (RED, "Skill compatibility", "Не каждая техника может выполнить\nлюбую заявку — проверка типа работ"),
    (PURPLE, "SLA + Priority", "high=+2ч, medium=+5ч, low=+12ч дедлайн\nШтраф за нарушение SLA в скоринге"),
]
for i, (color, title, body) in enumerate(features):
    col, row = i % 3, i // 3
    lx, ty = 0.3 + col * 4.25, 1.5 + row * 2.85
    rect(s, lx, ty, 4.0, 2.6, WHITE)
    rect(s, lx, ty, 4.0, 0.45, color)
    txt(s, title, lx + 0.15, ty + 0.07, 3.7, 0.35, size=13, bold=True, color=WHITE)
    txt(s, body, lx + 0.15, ty + 0.55, 3.7, 1.8, size=11, color=TEXT)
footer(s, 8)

# ═══ SLIDE 9 — ROADMAP ═══
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "Roadmap — от MVP к продакшену")
items = [
    (BLUE, "🔐", "Keycloak SSO", "Авторизация диспетчера, роли,\nавтоматическое определение смены"),
    (GREEN, "📡", "Real-time GPS", "WebSocket обновление позиций\nтехники каждые 30 сек из Wialon"),
    (RED, "🖥️", "GPU / локальный LLM", "vLLM для обоснований на русском\nбез отправки данных в облако"),
    (AMBER, "⚡", "Масштабирование", "Redis кеш матриц, Celery очередь\nVRP-задач, Gunicorn + N workers"),
]
for i, (color, icon, title, body) in enumerate(items):
    col, row = i % 2, i // 2
    lx, ty = 0.3 + col * 6.52, 1.5 + row * 2.8
    rect(s, lx, ty, 6.22, 2.55, WHITE)
    rect(s, lx, ty, 0.1, 2.55, color)
    txt(s, icon, lx + 0.2, ty + 0.15, 0.5, 0.5, size=24)
    txt(s, title, lx + 0.75, ty + 0.18, 5.2, 0.35, size=14, bold=True, color=color)
    txt(s, body, lx + 0.2, ty + 0.65, 5.8, 1.7, size=12, color=TEXT)
footer(s, 9)

# ═══ SLIDE 10 — THANK YOU ═══
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, BLUE)
rect(s, 0, 5.9, 13.33, 0.12, AMBER)
txt(s, "Спасибо за внимание", 0.5, 0.6, 12.33, 0.9, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
summary = [
    ("🎯", "35%", "Алгоритм", "OR-Tools VRP · 79% экономия · Dijkstra · скоринг"),
    ("⚙️", "25%", "Техника", "FastAPI · Docker · async · ThreadPool · PostgreSQL"),
    ("💡", "20%", "Инновации", "Gemini LLM · SA · SSSP O(N) · KDTree"),
    ("📊", "15%", "Применимость", "117 машин · 3293 скважины · Streamlit · Folium"),
    ("🔧", "10%", "Данные", "4624 узла · 38062 рёбер · snap 30ms · VRP 5s"),
]
for i, (icon, weight, title, detail) in enumerate(summary):
    lx = 0.35 + (i % 3) * 4.2
    ty = 1.65 + (i // 3) * 2.1
    rect(s, lx, ty, 3.9, 1.85, rgb(0x0D, 0x1B, 0x3E))
    rect(s, lx, ty, 3.9, 0.08, AMBER)
    txt(s, f"{icon}  {weight}", lx + 0.15, ty + 0.15, 3.6, 0.4, size=13, bold=True, color=AMBER)
    txt(s, title, lx + 0.15, ty + 0.56, 3.6, 0.35, size=12, bold=True, color=WHITE)
    txt(s, detail, lx + 0.15, ty + 0.96, 3.6, 0.7, size=10, color=rgb(0xAA, 0xCC, 0xFF))
txt(s, "Frontend: http://89.207.255.254:8503   ·   API: http://89.207.255.254:8003/docs",
    0.5, 6.9, 12.33, 0.35, size=12, color=rgb(0x99, 0xBB, 0xEE), align=PP_ALIGN.CENTER, italic=True)
footer(s, 10)

# ═══ SAVE ═══
OUT = "presentation_ИС_УТО.pptx"
prs.save(OUT)
print(f"✅  Saved: {OUT}  ({len(prs.slides)} slides)")
